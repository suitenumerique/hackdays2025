import pptx
def extract_text_from_pptx(file_obj):
    text = ""
    try:
        presentation = pptx.Presentation(file_obj)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Erreur lors de l'extraction du PPTX: {e}")
    return text